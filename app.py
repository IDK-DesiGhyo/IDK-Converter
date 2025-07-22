import os
import io
import zipfile
import base64
from datetime import datetime
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from PIL import Image, ImageDraw, ImageFont
import fitz  # PyMuPDF
from docx import Document
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches
import math
import tempfile
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)  # Enable CORS for GitHub Pages frontend

# Configuration
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB
UPLOAD_FOLDER = 'temp_uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

class EnhancedConverter:
    def __init__(self):
        self.supported_input_formats = {
            'image': ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif', '.webp'],
            'pdf': ['.pdf']
        }
        self.supported_output_formats = ['jpg', 'png', 'pdf', 'docx', 'pptx']

    def compress_image(self, image, target_size, format='JPEG'):
        """Compress image to target file size"""
        quality = 95
        min_quality = 20

        while quality >= min_quality:
            img_bytes = io.BytesIO()
            
            # Handle PNG with transparency
            if format == 'PNG':
                image.save(img_bytes, format='PNG', optimize=True)
            else:
                # Convert to RGB for JPEG
                if image.mode in ('RGBA', 'LA', 'P'):
                    rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                    if image.mode == 'P':
                        image = image.convert('RGBA')
                    rgb_image.paste(image, mask=image.split()[-1] if len(image.split()) == 4 else None)
                    image = rgb_image
                image.save(img_bytes, format='JPEG', quality=quality, optimize=True)

            if img_bytes.tell() <= target_size or format == 'PNG':
                img_bytes.seek(0)
                return img_bytes.getvalue()

            quality -= 5

        # If still too large, reduce dimensions
        width, height = image.size
        while img_bytes.tell() > target_size:
            width = int(width * 0.9)
            height = int(height * 0.9)
            resized_image = image.resize((width, height), Image.Resampling.LANCZOS)

            img_bytes = io.BytesIO()
            if format == 'PNG':
                resized_image.save(img_bytes, format='PNG', optimize=True)
            else:
                if resized_image.mode in ('RGBA', 'LA', 'P'):
                    rgb_image = Image.new('RGB', resized_image.size, (255, 255, 255))
                    if resized_image.mode == 'P':
                        resized_image = resized_image.convert('RGBA')
                    rgb_image.paste(resized_image, mask=resized_image.split()[-1] if len(resized_image.split()) == 4 else None)
                    resized_image = rgb_image
                resized_image.save(img_bytes, format='JPEG', quality=min_quality, optimize=True)

            if width < 100 or height < 100:  # Minimum size threshold
                break

        img_bytes.seek(0)
        return img_bytes.getvalue()

    def pdf_to_images(self, pdf_data):
        """Convert PDF to list of PIL Images"""
        images = []
        try:
            doc = fitz.open(stream=pdf_data, filetype="pdf")
            
            for page_num in range(len(doc)):
                page = doc[page_num]
                mat = fitz.Matrix(2, 2)  # 2x zoom for better quality
                pix = page.get_pixmap(matrix=mat)
                img_data = pix.tobytes("ppm")
                pil_image = Image.open(io.BytesIO(img_data))
                images.append(pil_image)
                pix = None  # Free memory
                
            doc.close()
        except Exception as e:
            logger.error(f"Error converting PDF: {str(e)}")
            raise e
        
        return images

    def create_collage(self, images, layout='grid', spacing=10, background_color='white'):
        """Create a collage from multiple images"""
        if not images:
            raise ValueError("No images provided for collage")

        # Calculate grid dimensions
        num_images = len(images)
        if layout == 'grid':
            cols = math.ceil(math.sqrt(num_images))
            rows = math.ceil(num_images / cols)
        elif layout == 'horizontal':
            cols = num_images
            rows = 1
        elif layout == 'vertical':
            cols = 1
            rows = num_images
        else:
            cols = math.ceil(math.sqrt(num_images))
            rows = math.ceil(num_images / cols)

        # Find the maximum dimensions for uniform sizing
        max_width = max(img.width for img in images)
        max_height = max(img.height for img in images)
        
        # Resize images to uniform size while maintaining aspect ratio
        uniform_images = []
        cell_width = max_width
        cell_height = max_height
        
        for img in images:
            # Calculate scaling to fit in cell while maintaining aspect ratio
            scale = min(cell_width / img.width, cell_height / img.height)
            new_width = int(img.width * scale)
            new_height = int(img.height * scale)
            
            resized_img = img.resize((new_width, new_height), Image.Resampling.LANCZOS)
            uniform_images.append(resized_img)

        # Calculate collage dimensions
        collage_width = cols * cell_width + (cols - 1) * spacing
        collage_height = rows * cell_height + (rows - 1) * spacing

        # Create collage
        collage = Image.new('RGB', (collage_width, collage_height), background_color)

        for i, img in enumerate(uniform_images):
            row = i // cols
            col = i % cols
            
            # Calculate position (center the image in its cell)
            x = col * (cell_width + spacing) + (cell_width - img.width) // 2
            y = row * (cell_height + spacing) + (cell_height - img.height) // 2
            
            collage.paste(img, (x, y))

        return collage

    def images_to_pdf(self, images):
        """Convert images to PDF"""
        pdf_buffer = io.BytesIO()
        
        # Convert all images to RGB
        rgb_images = []
        for img in images:
            if img.mode != 'RGB':
                rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'RGBA':
                    rgb_img.paste(img, mask=img.split()[-1])
                else:
                    rgb_img.paste(img)
                rgb_images.append(rgb_img)
            else:
                rgb_images.append(img)
        
        if rgb_images:
            rgb_images[0].save(pdf_buffer, format='PDF', save_all=True, append_images=rgb_images[1:])
        
        pdf_buffer.seek(0)
        return pdf_buffer.getvalue()

    def images_to_docx(self, images, title="Converted Images"):
        """Convert images to DOCX document"""
        doc = Document()
        doc.add_heading(title, 0)
        
        for i, img in enumerate(images):
            # Add image to document
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            doc.add_heading(f'Image {i + 1}', level=1)
            doc.add_picture(img_buffer, width=Inches(6))
            doc.add_page_break()
        
        # Save to buffer
        docx_buffer = io.BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)
        return docx_buffer.getvalue()

    def images_to_pptx(self, images, title="Converted Images"):
        """Convert images to PowerPoint presentation"""
        prs = Presentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        
        # Add images to slides
        for i, img in enumerate(images):
            slide_layout = prs.slide_layouts[6]  # Blank slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Add image
            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)
            
            # Calculate image position and size
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # Maintain aspect ratio
            img_ratio = img.width / img.height
            slide_ratio = slide_width / slide_height
            
            if img_ratio > slide_ratio:
                # Image is wider than slide ratio
                width = slide_width * 0.9
                height = width / img_ratio
            else:
                # Image is taller than slide ratio
                height = slide_height * 0.9
                width = height * img_ratio
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            slide.shapes.add_picture(img_buffer, left, top, width, height)
        
        # Save to buffer
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        return pptx_buffer.getvalue()

converter = EnhancedConverter()

@app.route('/', methods=['GET'])
def health_check():
    return jsonify({
        'status': 'healthy',
        'message': 'Enhanced PDF & Image Converter API',
        'supported_input': converter.supported_input_formats,
        'supported_output': converter.supported_output_formats
    })

@app.route('/convert', methods=['POST'])
def convert_files():
    try:
        data = request.get_json()
        
        if not data or 'files' not in data:
            return jsonify({'error': 'No files provided'}), 400
        
        files_data = data['files']
        output_format = data.get('output_format', 'jpg').lower()
        max_file_size = data.get('max_file_size', 500) * 1024  # Convert KB to bytes
        collage_settings = data.get('collage', {})
        
        if output_format not in converter.supported_output_formats:
            return jsonify({'error': f'Unsupported output format: {output_format}'}), 400
        
        all_images = []
        
        # Process each file
        for file_data in files_data:
            try:
                # Decode base64 file data
                file_content = base64.b64decode(file_data['content'])
                filename = file_data['name']
                file_ext = os.path.splitext(filename)[1].lower()
                
                if file_ext == '.pdf':
                    # Convert PDF to images
                    images = converter.pdf_to_images(file_content)
                    all_images.extend(images)
                elif file_ext in converter.supported_input_formats['image']:
                    # Load image
                    img = Image.open(io.BytesIO(file_content))
                    all_images.append(img)
                else:
                    logger.warning(f"Unsupported file format: {file_ext}")
                    continue
                    
            except Exception as e:
                logger.error(f"Error processing file {filename}: {str(e)}")
                continue
        
        if not all_images:
            return jsonify({'error': 'No valid images to process'}), 400
        
        # Create collage if requested and multiple images
        if collage_settings.get('enabled', False) and len(all_images) > 1:
            layout = collage_settings.get('layout', 'grid')
            spacing = collage_settings.get('spacing', 10)
            bg_color = collage_settings.get('background_color', 'white')
            
            collage_image = converter.create_collage(all_images, layout, spacing, bg_color)
            all_images = [collage_image]  # Use collage as the only image
        
        # Convert to requested format
        converted_files = []
        
        if output_format == 'jpg':
            for i, img in enumerate(all_images):
                compressed_data = converter.compress_image(img, max_file_size, 'JPEG')
                converted_files.append({
                    'name': f'converted_{i + 1}.jpg',
                    'content': base64.b64encode(compressed_data).decode('utf-8'),
                    'size': len(compressed_data)
                })
        
        elif output_format == 'png':
            for i, img in enumerate(all_images):
                compressed_data = converter.compress_image(img, max_file_size, 'PNG')
                converted_files.append({
                    'name': f'converted_{i + 1}.png',
                    'content': base64.b64encode(compressed_data).decode('utf-8'),
                    'size': len(compressed_data)
                })
        
        elif output_format == 'pdf':
            pdf_data = converter.images_to_pdf(all_images)
            converted_files.append({
                'name': 'converted_images.pdf',
                'content': base64.b64encode(pdf_data).decode('utf-8'),
                'size': len(pdf_data)
            })
        
        elif output_format == 'docx':
            docx_data = converter.images_to_docx(all_images, "Converted Images")
            converted_files.append({
                'name': 'converted_images.docx',
                'content': base64.b64encode(docx_data).decode('utf-8'),
                'size': len(docx_data)
            })
        
        elif output_format == 'pptx':
            pptx_data = converter.images_to_pptx(all_images, "Converted Images")
            converted_files.append({
                'name': 'converted_images.pptx',
                'content': base64.b64encode(pptx_data).decode('utf-8'),
                'size': len(pptx_data)
            })
        
        return jsonify({
            'success': True,
            'files': converted_files,
            'total_processed': len(all_images),
            'output_format': output_format
        })
        
    except Exception as e:
        logger.error(f"Conversion error: {str(e)}")
        return jsonify({'error': f'Conversion failed: {str(e)}'}), 500

@app.route('/collage', methods=['POST'])
def create_collage_only():
    try:
        data = request.get_json()
        
        if not data or 'files' not in data:
            return jsonify({'error': 'No files provided'}), 400
        
        files_data = data['files']
        layout = data.get('layout', 'grid')
        spacing = data.get('spacing', 10)
        bg_color = data.get('background_color', 'white')
        output_format = data.get('output_format', 'jpg').lower()
        max_file_size = data.get('max_file_size', 500) * 1024
        
        images = []
        
        # Process each file
        for file_data in files_data:
            try:
                file_content = base64.b64decode(file_data['content'])
                filename = file_data['name']
                file_ext = os.path.splitext(filename)[1].lower()
                
                if file_ext == '.pdf':
                    pdf_images = converter.pdf_to_images(file_content)
                    images.extend(pdf_images)
                elif file_ext in converter.supported_input_formats['image']:
                    img = Image.open(io.BytesIO(file_content))
                    images.append(img)
                    
            except Exception as e:
                logger.error(f"Error processing file {filename}: {str(e)}")
                continue
        
        if len(images) < 2:
            return jsonify({'error': 'At least 2 images required for collage'}), 400
        
        # Create collage
        collage_image = converter.create_collage(images, layout, spacing, bg_color)
        
        # Convert to requested format
        if output_format == 'png':
            compressed_data = converter.compress_image(collage_image, max_file_size, 'PNG')
            file_ext = 'png'
        else:
            compressed_data = converter.compress_image(collage_image, max_file_size, 'JPEG')
            file_ext = 'jpg'
        
        return jsonify({
            'success': True,
            'collage': {
                'name': f'collage.{file_ext}',
                'content': base64.b64encode(compressed_data).decode('utf-8'),
                'size': len(compressed_data)
            },
            'images_used': len(images)
        })
        
    except Exception as e:
        logger.error(f"Collage creation error: {str(e)}")
        return jsonify({'error': f'Collage creation failed: {str(e)}'}), 500

if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=False)